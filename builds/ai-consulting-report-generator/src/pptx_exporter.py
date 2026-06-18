"""PPTX exporter — Build 5 Phase 8.

Generates a 15-slide executive PowerPoint deck from generated Build 5 outputs.
Uses python-pptx only. No external AI calls. No real client data.
"""

import io
import os
import re
from datetime import date

_NAVY = (26, 39, 68)       # #1a2744
_BLUE = (37, 99, 235)      # #2563eb
_WHITE = (255, 255, 255)
_OFF_WHITE = (248, 250, 252)
_DARK = (30, 41, 59)       # #1e293b
_MUTED = (100, 116, 139)   # #64748b
_LIGHT = (226, 232, 240)   # #e2e8f0

_RESPONSIBLE_USE_BULLETS = [
    "All outputs generated from synthetic/demo audit data only.",
    "Must not be used with real client, learner, HR, or personal data.",
    "Does not provide legal, safeguarding, HR, compliance, or financial advice.",
    "Human review required before any real-world use.",
    "This is a consulting support artefact, not a professional advisory output.",
]


# ── Filename and text helpers ──────────────────────────────────────────────────


def create_safe_pptx_filename(title: str) -> str:
    """Return a safe dated PPTX filename."""
    slug = re.sub(r"[^\w\s-]", "", (title or "report").lower()).strip()
    slug = re.sub(r"[\s_-]+", "-", slug)[:60]
    return f"{date.today()}-{slug}.pptx"


def format_slide_text(text: str, max_chars: int = 700) -> str:
    """Strip Markdown and truncate text for slide use."""
    if not text:
        return ""
    t = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"\*(.+?)\*", r"\1", t)
    t = re.sub(r"---+", "", t)
    t = " ".join(t.split())
    if len(t) > max_chars:
        t = t[:max_chars].rstrip() + "…"
    return t.strip()


def extract_slide_bullets(items, max_items: int = 6) -> list:
    """Extract up to max_items bullets from a list or string."""
    if not items:
        return []
    if isinstance(items, str):
        lines = [l.strip().lstrip("- •").strip() for l in items.splitlines() if l.strip()]
        lines = [l for l in lines if l]
        return lines[:max_items]
    return [str(i).strip().lstrip("- •").strip() for i in items[:max_items]]


# ── Slide helpers ──────────────────────────────────────────────────────────────


def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _inches(n):
    from pptx.util import Inches
    return Inches(n)


def _pt(n):
    from pptx.util import Pt
    return Pt(n)


def _add_rect(slide, left, top, width, height, colour_rgb):
    """Add a filled rectangle to a slide."""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(*colour_rgb)
    shape.line.fill.background()
    return shape


def _add_textbox(
    slide,
    left, top, width, height,
    text,
    font_size=12,
    bold=False,
    colour_rgb=None,
    word_wrap=True,
    align_center=False,
):
    """Add a textbox with simple text."""
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = _pt(font_size)
    p.font.bold = bold
    if colour_rgb:
        p.font.color.rgb = _rgb(*colour_rgb)
    if align_center:
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    return txBox


def _add_bullet_textbox(
    slide,
    left, top, width, height,
    bullets,
    font_size=11,
    colour_rgb=None,
):
    """Add a textbox with bullet list paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    colour = colour_rgb or _DARK
    for i, bullet in enumerate(bullets or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = _pt(font_size)
        p.font.color.rgb = _rgb(*colour)
        p.space_after = _pt(4)
    return txBox


# ── Slide type builders ────────────────────────────────────────────────────────


def add_title_slide(
    prs,
    title: str,
    subtitle: str = "",
    footer: str = "",
) -> None:
    """Add a dark navy title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    W = prs.slide_width
    H = prs.slide_height

    # Navy background
    _add_rect(slide, 0, 0, W, H, _NAVY)

    # Top accent bar
    _add_rect(slide, 0, 0, W, _inches(0.08), _BLUE)

    # Title
    _add_textbox(
        slide,
        _inches(0.6), _inches(1.2), _inches(12), _inches(2.5),
        title,
        font_size=36, bold=True, colour_rgb=_WHITE,
    )

    # Subtitle
    if subtitle:
        _add_textbox(
            slide,
            _inches(0.6), _inches(3.8), _inches(12), _inches(1.2),
            subtitle,
            font_size=16, colour_rgb=(203, 213, 225),
        )

    # Footer
    if footer:
        _add_textbox(
            slide,
            _inches(0.6), _inches(6.5), _inches(12), _inches(0.6),
            footer,
            font_size=9, colour_rgb=(148, 163, 184),
        )


def add_bullet_slide(
    prs,
    title: str,
    bullets: list,
    speaker_notes: str = "",
) -> None:
    """Add a white-background slide with title and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    W = prs.slide_width

    # Navy title bar
    _add_rect(slide, 0, 0, W, _inches(1.1), _NAVY)

    # Title text
    _add_textbox(
        slide,
        _inches(0.4), _inches(0.15), _inches(12), _inches(0.8),
        title,
        font_size=22, bold=True, colour_rgb=_WHITE,
    )

    # Bullets
    if bullets:
        _add_bullet_textbox(
            slide,
            _inches(0.5), _inches(1.3), _inches(12.2), _inches(5.8),
            bullets,
            font_size=13,
        )

    # Speaker notes
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


def add_two_column_slide(
    prs,
    title: str,
    left_title: str,
    left_bullets: list,
    right_title: str,
    right_bullets: list,
) -> None:
    """Add a two-column bullet slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    W = prs.slide_width

    # Navy title bar
    _add_rect(slide, 0, 0, W, _inches(1.1), _NAVY)
    _add_textbox(
        slide,
        _inches(0.4), _inches(0.15), _inches(12), _inches(0.8),
        title, font_size=22, bold=True, colour_rgb=_WHITE,
    )

    # Left column label
    if left_title:
        _add_textbox(
            slide,
            _inches(0.4), _inches(1.3), _inches(6), _inches(0.45),
            left_title, font_size=11, bold=True, colour_rgb=_NAVY,
        )
    _add_bullet_textbox(
        slide,
        _inches(0.4), _inches(1.8), _inches(6), _inches(5.2),
        left_bullets, font_size=11,
    )

    # Right column label
    if right_title:
        _add_textbox(
            slide,
            _inches(6.8), _inches(1.3), _inches(6), _inches(0.45),
            right_title, font_size=11, bold=True, colour_rgb=_NAVY,
        )
    _add_bullet_textbox(
        slide,
        _inches(6.8), _inches(1.8), _inches(6), _inches(5.2),
        right_bullets, font_size=11,
    )


def add_chart_slide(
    prs,
    title: str,
    image_path: str,
    bullets: list | None = None,
) -> None:
    """Add a chart image slide. Falls back to text-only if image is missing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    W = prs.slide_width
    H = prs.slide_height

    _add_rect(slide, 0, 0, W, _inches(1.1), _NAVY)
    _add_textbox(
        slide,
        _inches(0.4), _inches(0.15), _inches(12), _inches(0.8),
        title, font_size=22, bold=True, colour_rgb=_WHITE,
    )

    image_added = False
    if image_path and os.path.exists(image_path):
        try:
            if bullets:
                # Chart on left, bullets on right
                slide.shapes.add_picture(
                    image_path,
                    _inches(0.3), _inches(1.2),
                    _inches(8.5), _inches(5.8),
                )
                _add_bullet_textbox(
                    slide,
                    _inches(9.0), _inches(1.5), _inches(4), _inches(5.5),
                    bullets, font_size=10,
                )
            else:
                # Chart centred
                slide.shapes.add_picture(
                    image_path,
                    _inches(0.5), _inches(1.2),
                    _inches(12.3), _inches(5.8),
                )
            image_added = True
        except Exception:
            pass

    if not image_added and bullets:
        _add_bullet_textbox(
            slide,
            _inches(0.5), _inches(1.3), _inches(12.3), _inches(5.8),
            bullets, font_size=12,
        )


def add_responsible_use_slide(prs) -> None:
    """Add the responsible-use boundaries slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    W = prs.slide_width
    # Amber accent to signal importance
    _add_rect(slide, 0, 0, W, _inches(1.1), (180, 83, 9))
    _add_textbox(
        slide,
        _inches(0.4), _inches(0.15), _inches(12), _inches(0.8),
        "Responsible-Use Boundaries",
        font_size=22, bold=True, colour_rgb=_WHITE,
    )
    _add_bullet_textbox(
        slide,
        _inches(0.5), _inches(1.3), _inches(12.3), _inches(5.8),
        _RESPONSIBLE_USE_BULLETS,
        font_size=13,
    )


# ── Full PPTX export ───────────────────────────────────────────────────────────


def export_client_report_to_pptx_bytes(
    export_data: dict,
    analytics: dict | None = None,
    chart_paths: dict | None = None,
) -> bytes:
    """Generate the executive PPTX deck. Returns bytes or b'' on failure."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        ed = export_data or {}
        an = analytics or {}
        cp = chart_paths or {}

        org = ed.get("organisation_name") or "Unnamed organisation"
        generated = ed.get("generated_date") or str(date.today())

        audit = ed.get("audit_data") or {}
        profile = audit.get("organisation_profile") or {}
        org_type = profile.get("organisation_type", "")
        sector = profile.get("sector", "")
        staff = profile.get("staff_count", "")
        current_ai = profile.get("current_ai_use", "")

        rs_data = ed.get("readiness_summary") or {}
        opp_summary = ed.get("opportunity_summary") or {}
        opp_portfolio = ed.get("opportunity_portfolio") or {}
        rm_summary = ed.get("implementation_roadmap_summary") or {}
        rr_summary = ed.get("risk_register_summary") or {}
        report_secs = ed.get("report_sections") or {}
        sections = report_secs.get("sections") or {} if isinstance(report_secs, dict) else {}

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        # Ensure blank layout exists
        while len(prs.slide_layouts) < 7:
            prs.slide_layouts._sldLayoutLst.append(
                prs.slide_layouts[0]._element
            )

        # ── Slide 1: Title ─────────────────────────────────────────────────────
        add_title_slide(
            prs,
            title="AI Readiness and\nResponsible AI Adoption Report",
            subtitle=org,
            footer=(
                f"Generated: {generated}  |  "
                "Production-style prototype — synthetic data only — human review required"
            ),
        )

        # ── Slide 2: Organisation Context ──────────────────────────────────────
        context_bullets = []
        if org_type:
            context_bullets.append(f"Organisation type: {org_type}")
        if sector:
            context_bullets.append(f"Sector: {sector}")
        if staff:
            context_bullets.append(f"Staff: {staff}")
        if current_ai:
            # Truncate long AI use text
            context_bullets.append(f"Current AI use: {format_slide_text(current_ai, 200)}")
        deps = profile.get("departments") or []
        if deps:
            context_bullets.append("Departments: " + ", ".join(deps[:4]))
        goals = profile.get("main_business_goals") or []
        context_bullets.extend(extract_slide_bullets(goals, 3))
        add_bullet_slide(
            prs,
            title="Organisation Context",
            bullets=context_bullets or ["No organisation profile available."],
            speaker_notes="Describe the organisation and its current AI context.",
        )

        # ── Slide 3: AI Readiness Summary ──────────────────────────────────────
        readiness_bullets = []
        if rs_data:
            overall = rs_data.get("overall_score", "")
            level = rs_data.get("overall_level", "")
            if overall and level:
                readiness_bullets.append(f"Overall AI readiness score: {int(overall)}/100 — {level}")
            ranked = rs_data.get("ranked_categories") or []
            if ranked:
                readiness_bullets.append(f"Strongest area: {ranked[0].get('label', '')} ({int(ranked[0].get('score', 0))}/100)")
                readiness_bullets.append(f"Weakest area: {ranked[-1].get('label', '')} ({int(ranked[-1].get('score', 0))}/100)")
            interp = rs_data.get("strategic_interpretation", "")
            if interp:
                readiness_bullets.append(format_slide_text(interp, 250))
        else:
            readiness_bullets.append("Readiness Summary not yet generated.")
            readiness_bullets.append("Run the Readiness Summary page to generate this data.")
        add_bullet_slide(
            prs,
            title="AI Readiness Summary",
            bullets=readiness_bullets,
            speaker_notes="Walk through the overall readiness score and strongest/weakest areas.",
        )

        # ── Slide 4: Readiness Score Chart ─────────────────────────────────────
        readiness_score_labels = []
        score_breakdown = an.get("readiness_score_breakdown") or {}
        for cat, score in score_breakdown.items():
            readiness_score_labels.append(f"{cat}: {score}/100")
        add_chart_slide(
            prs,
            title="AI Readiness Scores by Category",
            image_path=cp.get("readiness_scores", ""),
            bullets=readiness_score_labels[:6] or None,
        )

        # ── Slide 5: Key Findings ──────────────────────────────────────────────
        kf_sec = sections.get("key_findings") or {}
        kf_bullets = extract_slide_bullets(kf_sec.get("key_points") or [], 6)
        if not kf_bullets:
            audit_d = ed.get("audit_data") or {}
            kf_bullets = [
                f"{len(audit_d.get('risk_findings') or [])} AI risks identified — review before scaling.",
                f"{len(audit_d.get('workflow_findings') or [])} workflow opportunities for AI pilots.",
                "Informal AI use already occurring — policy needed urgently.",
                "Sensitive data boundaries must be defined before AI tools are used.",
                f"{len(audit_d.get('training_needs') or [])} training needs identified.",
                "Controlled pilots recommended over broad rollout.",
            ]
        add_bullet_slide(
            prs, title="Key Findings", bullets=kf_bullets,
            speaker_notes="Summarise the most important findings from the audit.",
        )

        # ── Slide 6: Risk Summary ──────────────────────────────────────────────
        risk_bullets = []
        if rr_summary:
            total = rr_summary.get("total_risks", 0)
            high = rr_summary.get("high_risks", 0)
            critical = rr_summary.get("critical_risks", 0)
            risk_bullets.append(f"Total risks identified: {total}")
            risk_bullets.append(f"Critical risks: {critical}  |  High risks: {high}")
            position = rr_summary.get("overall_risk_position", "")
            if position:
                risk_bullets.append(format_slide_text(position, 200))
            focus = rr_summary.get("recommended_focus") or []
            risk_bullets.extend(extract_slide_bullets(focus, 3))
        else:
            risk_bullets.append("Risk Register not yet generated.")
            risk_bullets.append("Run the Risk Register page to generate this data.")
        add_bullet_slide(
            prs, title="AI Risk Summary",
            bullets=risk_bullets,
            speaker_notes="Walk through the risk position and priority controls.",
        )

        # ── Slide 7: Risk Chart ────────────────────────────────────────────────
        risk_counts = an.get("risk_level_counts") or {}
        risk_labels = [f"{level}: {count}" for level, count in risk_counts.items() if count > 0]
        add_chart_slide(
            prs,
            title="AI Risk Distribution",
            image_path=cp.get("risk_levels", ""),
            bullets=risk_labels or None,
        )

        # ── Slide 8: Opportunity and Pilot Recommendations ─────────────────────
        opp_bullets = []
        if opp_summary:
            opp_bullets.append(
                f"Total AI opportunities: {opp_summary.get('total_opportunities', 0)}"
            )
            opp_bullets.append(
                f"Recommended pilots: {opp_summary.get('total_pilots', 0)}"
            )
            first_pilot = opp_summary.get("recommended_first_pilot_name", "")
            if first_pilot:
                opp_bullets.append(f"Recommended first pilot: {first_pilot}")
            position = opp_summary.get("overall_opportunity_position", "")
            if position:
                opp_bullets.append(format_slide_text(position, 200))
            focus = opp_summary.get("recommended_focus") or []
            opp_bullets.extend(extract_slide_bullets(focus, 2))
        if opp_portfolio:
            sequence = opp_portfolio.get("recommended_pilot_sequence") or []
            for item in sequence[:3]:
                opp_bullets.append(
                    f"Pilot {item.get('position', '')}: {item.get('pilot_name', '')} "
                    f"— {item.get('suggested_timeline', '')}"
                )
        if not opp_bullets:
            opp_bullets.append("Opportunity Portfolio not yet generated.")
        add_bullet_slide(
            prs, title="AI Opportunity and Pilot Recommendations",
            bullets=opp_bullets,
            speaker_notes="Present the recommended pilot sequence and first pilot.",
        )

        # ── Slide 9: Opportunity Chart ─────────────────────────────────────────
        opp_counts = an.get("opportunity_priority_counts") or {}
        opp_labels = [f"{p}: {c}" for p, c in opp_counts.items() if c > 0]
        add_chart_slide(
            prs,
            title="AI Opportunity Priority Distribution",
            image_path=cp.get("opportunity_priorities", ""),
            bullets=opp_labels or None,
        )

        # ── Slide 10: 30/60/90-Day Roadmap ────────────────────────────────────
        roadmap_bullets = []
        if rm_summary:
            roadmap_bullets.append(
                f"Total actions: {rm_summary.get('total_actions', 0)}"
            )
            roadmap_bullets.append(
                f"First 30 days: {rm_summary.get('day_30_actions', 0)} actions — "
                "governance, data controls, risk management"
            )
            roadmap_bullets.append(
                f"Days 31–60: {rm_summary.get('day_60_actions', 0)} actions — "
                "pilot preparation and controlled delivery"
            )
            roadmap_bullets.append(
                f"Days 61–90: {rm_summary.get('day_90_actions', 0)} actions — "
                "review, refine, and scale decision"
            )
            pilot = rm_summary.get("recommended_first_pilot", "")
            if pilot:
                roadmap_bullets.append(f"Recommended first pilot: {pilot}")
            position = rm_summary.get("overall_roadmap_position", "")
            if position:
                roadmap_bullets.append(format_slide_text(position, 200))
        else:
            roadmap_bullets.append("Implementation Roadmap not yet generated.")
        add_bullet_slide(
            prs, title="30/60/90-Day Implementation Roadmap",
            bullets=roadmap_bullets,
            speaker_notes="Walk through the three roadmap phases.",
        )

        # ── Slide 11: Roadmap Chart ────────────────────────────────────────────
        rm_counts = an.get("roadmap_action_counts") or {}
        rm_labels = [f"{phase}: {count} actions" for phase, count in rm_counts.items() if count > 0]
        add_chart_slide(
            prs,
            title="Roadmap Actions by Phase",
            image_path=cp.get("roadmap_actions", ""),
            bullets=rm_labels or None,
        )

        # ── Slide 12: Governance Recommendations ──────────────────────────────
        gov_sec = sections.get("governance_recommendations") or {}
        gov_bullets = extract_slide_bullets(gov_sec.get("recommendations") or [], 6)
        if not gov_bullets:
            gov_bullets = [
                "Define approved and prohibited AI use cases.",
                "Clarify data boundaries — specify what may never enter AI tools.",
                "Clarify safeguarding escalation routes.",
                "Assign AI governance ownership to a named responsible owner.",
                "Require human review before any AI-generated output is used.",
                "Run staff training before scaling AI use.",
            ]
        add_bullet_slide(
            prs, title="Governance Recommendations",
            bullets=gov_bullets,
            speaker_notes="Present the governance structure needed before scaling AI.",
        )

        # ── Slide 13: Immediate Next Steps ─────────────────────────────────────
        next_sec = sections.get("immediate_next_steps") or {}
        next_bullets = extract_slide_bullets(next_sec.get("key_points") or [], 6)
        if not next_bullets:
            next_bullets = [
                "Confirm responsible owners for AI governance.",
                "Review and act on the risk register.",
                "Approve the scope of the first pilot.",
                "Define success measures and review date.",
                "Prepare and deliver staff training.",
                "Run the controlled pilot with human review.",
            ]
        add_bullet_slide(
            prs, title="Immediate Next Steps",
            bullets=next_bullets,
            speaker_notes="Agree next steps before the client leaves.",
        )

        # ── Slide 14: Responsible-Use Boundaries ───────────────────────────────
        add_responsible_use_slide(prs)

        # ── Slide 15: Recommended Direction ───────────────────────────────────
        direction_bullets = [
            "Strengthen governance and data controls before scaling AI use.",
            "Run one controlled pilot in the safest, lowest-complexity workflow.",
            "Deliver responsible AI use training to all staff before the pilot.",
            "Review pilot outcomes against success measures before any further scaling.",
            "Ensure human review at every step — AI outputs support humans, not replace them.",
            "This report is a starting point, not a final plan — consultant review required.",
        ]
        add_bullet_slide(
            prs, title="Recommended Direction",
            bullets=direction_bullets,
            speaker_notes=(
                "Close by reinforcing the responsible direction: controlled, evidence-based, "
                "human-reviewed AI adoption."
            ),
        )

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    except Exception:
        return b""
