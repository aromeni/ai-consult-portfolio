"""Tests for src/pptx_exporter.py — 15-slide Training Deck export."""

import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from src import pptx_exporter as ppt


# ── Shared fixtures ─────────────────────────────────────────────────────────────

MINIMAL_PACK: dict = {}

FULL_PACK: dict = {
    "scenario": {
        "organisation_name": "BrightPath Skills Training",
        "organisation_type": "Training provider",
        "sector": "Education",
        "staff_count": 8,
        "staff_roles": ["Tutors", "Administrators", "Team Leaders"],
        "main_concerns": ["learner data", "safeguarding", "hallucination"],
        "priority_topics": ["Learner data protection", "Hallucination risk", "Bias"],
        "current_ai_use": "Informal ChatGPT use for lesson planning and emails",
        "training_goal": "Use AI safely for lesson planning and admin",
    },
    "training_needs_assessment": {
        "topic_assessments": [
            {"title": "Learner Data Protection", "priority_level": "high"},
            {"title": "Safeguarding Boundaries", "priority_level": "high"},
            {"title": "Hallucination Awareness", "priority_level": "medium"},
        ],
        "recommended_learning_outcomes": [
            "Identify what information must not be entered into AI tools.",
            "Apply the human review checklist before using AI outputs.",
            "Recognise hallucination and verify AI claims independently.",
        ],
    },
    "workshop_plan": {
        "workshop_title": "Using AI Safely for Lesson Planning",
        "duration_minutes": 90,
        "delivery_mode": "In-person",
        "agenda": [
            {"time_range": "00:00-00:10", "section_title": "Welcome and Ground Rules"},
            {"time_range": "00:10-00:25", "section_title": "What is AI?"},
            {"time_range": "00:25-00:50", "section_title": "Safe Prompting Practice"},
        ],
        "responsible_use_messages": [
            "Never enter learner data into AI tools.",
            "AI outputs must be reviewed before use.",
        ],
    },
    "training_activities": [
        {
            "activity_title": "Safe vs Unsafe Prompt Sorting",
            "duration_minutes": 15,
            "activity_type": "sorting",
        },
        {
            "activity_title": "Spot the Hallucination",
            "duration_minutes": 10,
            "activity_type": "hallucination_review",
        },
    ],
    "staff_handout": {
        "allowed_ai_uses": [
            "Draft generic lesson plans without learner names.",
            "Rewrite professional emails for tone.",
        ],
        "prohibited_ai_uses": [
            "Enter learner names or personal data.",
            "Use AI to make safeguarding decisions.",
        ],
        "human_review_checklist": [
            "Check all facts before using AI output.",
            "Remove any personal data from prompts.",
        ],
        "escalation_guidance": [
            {"issue": "Safeguarding concern", "what_to_do": "Contact the DSL immediately."},
            {"issue": "Data breach", "what_to_do": "Report to DPO within 24 hours."},
        ],
    },
    "knowledge_check": {
        "multiple_choice_questions": [
            {"question": "Which prompt is safest?", "options": ["A", "B", "C", "D"], "correct": "B"},
        ],
        "scenario_questions": [
            {"question": "A tutor wants to paste a safeguarding concern into ChatGPT. What should they do?"},
        ],
        "reflection_questions": [
            {"question": "What is one type of information you must never enter into an AI tool?"},
        ],
        "answer_key": True,
        "pass_guidance": "8 out of 10 correct to pass.",
    },
}

FULL_ANALYTICS: dict = {
    "section_completion": {
        "Organisation Scenario": True,
        "Training Needs Assessment": True,
        "Workshop Plan": True,
        "Training Activities": False,
        "Facilitator Guide": False,
    },
    "activity_type_counts": {"sorting": 2, "hallucination_review": 1},
    "knowledge_check_topic_counts": {"learner_data": 3, "safeguarding": 2},
    "report_quality": {
        "sections_available": 3,
        "sections_total": 7,
        "completeness_pct": 42.8,
    },
}


def _prs():
    """Return a blank widescreen Presentation for testing slide builders."""
    p = Presentation()
    p.slide_width = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p


# ── create_safe_pptx_filename ───────────────────────────────────────────────────

class TestCreateSafePptxFilename:
    def test_returns_pptx_extension(self):
        assert ppt.create_safe_pptx_filename("Training Pack").endswith(".pptx")

    def test_slugifies_spaces(self):
        result = ppt.create_safe_pptx_filename("My Training Pack")
        assert " " not in result
        assert "my-training-pack" in result

    def test_slugifies_special_chars(self):
        result = ppt.create_safe_pptx_filename("BrightPath! Pack (2026)")
        assert "!" not in result
        assert "(" not in result

    def test_non_string_falls_back(self):
        result = ppt.create_safe_pptx_filename(None)
        assert result.endswith(".pptx")
        assert len(result) > 5

    def test_empty_string_falls_back(self):
        result = ppt.create_safe_pptx_filename("")
        assert result.endswith(".pptx")

    def test_max_slug_length(self):
        result = ppt.create_safe_pptx_filename("a" * 200)
        slug = result[: -len(".pptx")]
        assert len(slug) <= 60


# ── format_slide_text ───────────────────────────────────────────────────────────

class TestFormatSlideText:
    def test_strips_bold_markdown(self):
        result = ppt.format_slide_text("**bold**")
        assert "**" not in result
        assert "bold" in result

    def test_strips_italic_markdown(self):
        result = ppt.format_slide_text("*italic*")
        assert result.count("*") == 0

    def test_strips_code_markdown(self):
        result = ppt.format_slide_text("`code`")
        assert "`" not in result
        assert "code" in result

    def test_strips_heading_markdown(self):
        result = ppt.format_slide_text("## Section Heading")
        assert "#" not in result
        assert "Section Heading" in result

    def test_truncates_at_max_chars(self):
        long_text = "word " * 200
        result = ppt.format_slide_text(long_text, max_chars=50)
        assert len(result) <= 55

    def test_adds_ellipsis_on_truncation(self):
        result = ppt.format_slide_text("a " * 500, max_chars=700)
        assert result.endswith("…")

    def test_short_text_not_truncated(self):
        result = ppt.format_slide_text("Short text here.")
        assert "Short text here." in result

    def test_none_returns_empty_string(self):
        assert ppt.format_slide_text(None) == ""

    def test_integer_returns_empty_string(self):
        assert ppt.format_slide_text(42) == ""

    def test_default_max_chars_is_700(self):
        text = "word " * 200
        result = ppt.format_slide_text(text)
        assert len(result) <= 705


# ── extract_slide_bullets ───────────────────────────────────────────────────────

class TestExtractSlideBullets:
    def test_list_of_strings(self):
        result = ppt.extract_slide_bullets(["A", "B", "C"])
        assert result == ["A", "B", "C"]

    def test_respects_max_items(self):
        result = ppt.extract_slide_bullets(list("ABCDEFG"), max_items=3)
        assert len(result) == 3

    def test_default_max_is_six(self):
        result = ppt.extract_slide_bullets(list("ABCDEFGHIJ"))
        assert len(result) == 6

    def test_string_splits_by_newline(self):
        result = ppt.extract_slide_bullets("Line 1\nLine 2\nLine 3")
        assert len(result) == 3
        assert "Line 1" in result[0]

    def test_list_of_dicts_extracts_title(self):
        items = [{"title": "Topic A"}, {"title": "Topic B"}]
        result = ppt.extract_slide_bullets(items)
        assert "Topic A" in result[0]
        assert "Topic B" in result[1]

    def test_list_of_dicts_falls_back_to_text(self):
        items = [{"text": "Body text"}]
        result = ppt.extract_slide_bullets(items)
        assert "Body text" in result[0]

    def test_empty_list_returns_empty(self):
        assert ppt.extract_slide_bullets([]) == []

    def test_none_returns_empty(self):
        assert ppt.extract_slide_bullets(None) == []

    def test_truncates_long_items(self):
        long_item = "word " * 100
        result = ppt.extract_slide_bullets([long_item])
        assert len(result[0]) <= 130


# ── Slide builder functions ─────────────────────────────────────────────────────

class TestSlideBuilders:
    def test_add_title_slide_adds_one_slide(self):
        p = _prs()
        ppt.add_title_slide(p, "Title", "Subtitle")
        assert len(p.slides) == 1

    def test_add_title_slide_no_subtitle(self):
        p = _prs()
        ppt.add_title_slide(p, "Title Only")
        assert len(p.slides) == 1

    def test_add_title_slide_custom_footer(self):
        p = _prs()
        ppt.add_title_slide(p, "T", footer="Custom footer text")
        assert len(p.slides) == 1

    def test_add_bullet_slide_adds_one_slide(self):
        p = _prs()
        ppt.add_bullet_slide(p, "Bullets", ["A", "B", "C"])
        assert len(p.slides) == 1

    def test_add_bullet_slide_empty_bullets(self):
        p = _prs()
        ppt.add_bullet_slide(p, "Title", [])
        assert len(p.slides) == 1

    def test_add_bullet_slide_with_speaker_notes(self):
        p = _prs()
        ppt.add_bullet_slide(p, "Title", ["Bullet"], speaker_notes="Notes here.")
        assert len(p.slides) == 1

    def test_add_two_column_slide_adds_one_slide(self):
        p = _prs()
        ppt.add_two_column_slide(p, "Title", "Left", ["L1", "L2"], "Right", ["R1", "R2"])
        assert len(p.slides) == 1

    def test_add_two_column_slide_empty_columns(self):
        p = _prs()
        ppt.add_two_column_slide(p, "Title", "", [], "", [])
        assert len(p.slides) == 1

    def test_add_responsible_use_slide_adds_one_slide(self):
        p = _prs()
        ppt.add_responsible_use_slide(p)
        assert len(p.slides) == 1

    def test_multiple_slides_accumulate(self):
        p = _prs()
        ppt.add_title_slide(p, "T1")
        ppt.add_bullet_slide(p, "T2", ["B"])
        ppt.add_two_column_slide(p, "T3", "L", ["L1"], "R", ["R1"])
        ppt.add_responsible_use_slide(p)
        assert len(p.slides) == 4


# ── export_training_pack_to_pptx_bytes ─────────────────────────────────────────

class TestExportTrainingPackToPptxBytes:
    def test_returns_bytes(self):
        result = ppt.export_training_pack_to_pptx_bytes(FULL_PACK)
        assert isinstance(result, bytes)

    def test_pptx_bytes_non_empty(self):
        result = ppt.export_training_pack_to_pptx_bytes(FULL_PACK)
        assert len(result) > 5000

    def test_pptx_magic_bytes(self):
        result = ppt.export_training_pack_to_pptx_bytes(FULL_PACK)
        assert result.startswith(b"PK")

    def test_missing_analytics_does_not_crash(self):
        result = ppt.export_training_pack_to_pptx_bytes(FULL_PACK, analytics=None)
        assert result.startswith(b"PK")

    def test_empty_pack_data_does_not_crash(self):
        result = ppt.export_training_pack_to_pptx_bytes({})
        assert result.startswith(b"PK")

    def test_minimal_pack_data_does_not_crash(self):
        result = ppt.export_training_pack_to_pptx_bytes(MINIMAL_PACK)
        assert isinstance(result, bytes)

    def test_none_pack_data_does_not_crash(self):
        result = ppt.export_training_pack_to_pptx_bytes(None)
        assert result.startswith(b"PK")

    def test_with_full_analytics(self):
        result = ppt.export_training_pack_to_pptx_bytes(FULL_PACK, analytics=FULL_ANALYTICS)
        assert result.startswith(b"PK")

    def test_produces_15_slides(self):
        result = ppt.export_training_pack_to_pptx_bytes(FULL_PACK, analytics=FULL_ANALYTICS)
        p = Presentation(io.BytesIO(result))
        assert len(p.slides) == 15

    def test_produces_15_slides_minimal_data(self):
        result = ppt.export_training_pack_to_pptx_bytes(MINIMAL_PACK)
        p = Presentation(io.BytesIO(result))
        assert len(p.slides) == 15

    def test_missing_staff_handout_does_not_crash(self):
        pack = {"scenario": FULL_PACK["scenario"]}
        result = ppt.export_training_pack_to_pptx_bytes(pack)
        assert result.startswith(b"PK")

    def test_missing_knowledge_check_does_not_crash(self):
        pack = {k: v for k, v in FULL_PACK.items() if k != "knowledge_check"}
        result = ppt.export_training_pack_to_pptx_bytes(pack)
        assert result.startswith(b"PK")

    def test_missing_workshop_plan_does_not_crash(self):
        pack = {k: v for k, v in FULL_PACK.items() if k != "workshop_plan"}
        result = ppt.export_training_pack_to_pptx_bytes(pack)
        assert result.startswith(b"PK")

    def test_section_completion_analytics_renders(self):
        analytics = {"section_completion": {"Scenario": True, "Needs": False}}
        result = ppt.export_training_pack_to_pptx_bytes(FULL_PACK, analytics=analytics)
        assert result.startswith(b"PK")
