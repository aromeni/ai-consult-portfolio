"""PPTX export — AI Staff Training and Workshop Generator.

Self-contained 15-slide professional training deck builder. Converts Build 4
generated training pack data into a PowerPoint presentation for demo and
portfolio review.

No external AI API calls. Synthetic scenario content only.
Outputs require human review before use in a real training context.
"""

import io
import re
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Colour palette ──────────────────────────────────────────────────────────────

_NAVY  = RGBColor(0x1A, 0x27, 0x44)
_BLUE  = RGBColor(0x3B, 0x6C, 0xF7)
_DARK  = RGBColor(0x1E, 0x29, 0x3B)
_GREY  = RGBColor(0x64, 0x74, 0x8B)
_MUTED = RGBColor(0xBD, 0xC3, 0xC7)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GREEN = RGBColor(0x16, 0xA3, 0x4A)
_RED   = RGBColor(0xDC, 0x26, 0x26)
_DIVIDER = RGBColor(0xE2, 0xE8, 0xF0)

_SLIDE_W = 13.33  # inches — widescreen 16:9
_SLIDE_H = 7.5
_FOOTER  = "Synthetic scenario prototype. Human review required before use."


# ── Utility functions ───────────────────────────────────────────────────────────

def create_safe_pptx_filename(title: str) -> str:
    """Return a kebab-case .pptx filename from a title string."""
    if not isinstance(title, str):
        title = "training-pack"
    slug = re.sub(r"[^a-z0-9]+", "-", title.lower().strip()).strip("-")
    slug = slug[:60] or "training-pack"
    return f"{slug}.pptx"


def format_slide_text(text: str, max_chars: int = 700) -> str:
    """Strip Markdown formatting and truncate text for slide use."""
    if not isinstance(text, str):
        return ""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"^[-*]\s+", "", text, flags=re.MULTILINE)
    text = text.strip()
    if len(text) > max_chars:
        text = text[:max_chars].rsplit(" ", 1)[0] + "…"
    return text


def extract_slide_bullets(items, max_items: int = 6) -> list:
    """Extract a clean list of bullet strings from various input shapes.

    Handles: list of strings, list of dicts (extracts title/text/label/item),
    newline-delimited string, None/empty.
    """
    if not items:
        return []
    if isinstance(items, str):
        lines = [ln.strip() for ln in items.splitlines() if ln.strip()]
        return [format_slide_text(ln, 120) for ln in lines[:max_items]]
    if isinstance(items, (list, tuple)):
        result = []
        for item in items[:max_items]:
            if isinstance(item, dict):
                text = (
                    item.get("title") or item.get("text") or
                    item.get("label") or item.get("item") or str(item)
                )
            else:
                text = str(item)
            result.append(format_slide_text(text, 120))
        return result
    return []


# ── Low-level primitives ────────────────────────────────────────────────────────

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide, text: str, left: float, top: float, width: float, height: float,
         font_size: int = 16, bold: bool = False, color=None,
         align=PP_ALIGN.LEFT, wrap: bool = True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def _navy_bar(slide):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(_SLIDE_W), Inches(1.35)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _NAVY
    bar.line.fill.background()


def _footer(slide, text: str = _FOOTER):
    _box(slide, text, 0.3, _SLIDE_H - 0.45, _SLIDE_W - 0.6, 0.4,
         font_size=9, color=_GREY)


def _notes(slide, text: str):
    if not text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def _bullet_body(bullets: list, max_items: int = 6) -> str:
    return "\n".join(f"• {b}" for b in bullets[:max_items])


# ── Public slide builders ───────────────────────────────────────────────────────

def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str = "",
    footer: str = "",
) -> None:
    """Full navy title slide — used as the first slide or section divider."""
    slide = _blank(prs)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(_SLIDE_W), Inches(_SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = _NAVY
    bg.line.fill.background()

    _box(slide, title, 0.5, 1.8, _SLIDE_W - 1.0, 1.5,
         font_size=36, bold=True, color=_WHITE)
    if subtitle:
        _box(slide, subtitle, 0.5, 3.45, _SLIDE_W - 1.0, 0.9,
             font_size=20, color=_BLUE)
    notice = footer or _FOOTER
    _box(slide, notice, 0.5, _SLIDE_H - 0.7, _SLIDE_W - 1.0, 0.55,
         font_size=10, color=_MUTED)


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list,
    speaker_notes: str = "",
) -> None:
    """Standard bullet list slide with navy title bar. Max 6 bullets."""
    slide = _blank(prs)
    _navy_bar(slide)
    _box(slide, title, 0.35, 0.08, _SLIDE_W - 0.7, 1.15,
         font_size=26, bold=True, color=_WHITE)
    if bullets:
        _box(slide, _bullet_body(bullets, 6), 0.5, 1.55, _SLIDE_W - 1.0, _SLIDE_H - 2.3,
             font_size=16, color=_DARK)
    _footer(slide)
    _notes(slide, speaker_notes)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_bullets: list,
    right_title: str,
    right_bullets: list,
) -> None:
    """Two-column slide with navy title bar and a subtle vertical divider."""
    slide = _blank(prs)
    _navy_bar(slide)
    _box(slide, title, 0.35, 0.08, _SLIDE_W - 0.7, 1.15,
         font_size=26, bold=True, color=_WHITE)

    col_w = (_SLIDE_W - 1.2) / 2
    right_x = col_w + 0.7

    # Vertical divider
    div = slide.shapes.add_shape(
        1, Inches(col_w + 0.59), Inches(1.5), Inches(0.02), Inches(_SLIDE_H - 2.2)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = _DIVIDER
    div.line.fill.background()

    # Left column
    _box(slide, left_title, 0.4, 1.55, col_w, 0.45,
         font_size=14, bold=True, color=_NAVY)
    _box(slide, _bullet_body(left_bullets, 6), 0.4, 2.05, col_w, _SLIDE_H - 2.9,
         font_size=14, color=_DARK)

    # Right column
    _box(slide, right_title, right_x, 1.55, col_w, 0.45,
         font_size=14, bold=True, color=_NAVY)
    _box(slide, _bullet_body(right_bullets, 6), right_x, 2.05, col_w, _SLIDE_H - 2.9,
         font_size=14, color=_DARK)

    _footer(slide)


def add_responsible_use_slide(prs: Presentation) -> None:
    """Slide 14 — non-negotiable responsible-use boundaries with red left accent."""
    bullets = [
        "Use synthetic or approved scenarios only — no real learner, safeguarding, HR, or personal data.",
        "No external AI API calls — all generation runs locally and deterministically.",
        "Human review is required before using any output in a real training session.",
        "This pack does not constitute legal, safeguarding, HR, compliance, or professional advice.",
        "Training materials are support materials — not official organisational policy.",
        "Production use requires governance, DPIA, security, and responsible-owner approval.",
    ]
    slide = _blank(prs)
    _navy_bar(slide)
    _box(slide, "Responsible-Use Boundaries", 0.35, 0.08, _SLIDE_W - 0.7, 1.15,
         font_size=26, bold=True, color=_WHITE)
    # Red left accent bar
    accent = slide.shapes.add_shape(
        1, Inches(0.3), Inches(1.5), Inches(0.08), Inches(_SLIDE_H - 2.2)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _RED
    accent.line.fill.background()
    _box(slide, _bullet_body(bullets, 6), 0.55, 1.55, _SLIDE_W - 0.85, _SLIDE_H - 2.3,
         font_size=15, color=_DARK)
    _footer(slide)


# ── Main export ─────────────────────────────────────────────────────────────────

def export_training_pack_to_pptx_bytes(
    pack_data: dict,
    analytics: dict | None = None,
) -> bytes:
    """Generate a 15-slide training deck and return its bytes.

    Suitable for st.download_button(data=..., mime="application/vnd.openxmlformats-
    officedocument.presentationml.presentation").

    Missing pack sections produce a "not yet generated" placeholder — the export
    never crashes on incomplete data.
    """
    if not isinstance(pack_data, dict):
        pack_data = {}
    if analytics is None:
        analytics = {}

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)

    scenario        = pack_data.get("scenario") or {}
    org_name        = scenario.get("organisation_name", "Organisation")
    org_type        = scenario.get("organisation_type", "")
    sector          = scenario.get("sector", "")
    staff_count     = scenario.get("staff_count", "")
    roles           = scenario.get("staff_roles") or []
    concerns        = scenario.get("main_concerns") or []
    current_ai_use  = scenario.get("current_ai_use", "")
    priority_topics = scenario.get("priority_topics") or []
    today           = str(date.today())

    # ── Slide 1: Title ──────────────────────────────────────────────────────────
    sub_parts = [p for p in [org_type, sector, f"{staff_count} staff" if staff_count else ""] if p]
    subtitle = f"{org_name}  ·  {'  |  '.join(sub_parts)}" if sub_parts else org_name
    add_title_slide(
        prs,
        title="Responsible AI Staff Training Pack",
        subtitle=subtitle,
        footer=f"Generated: {today}  ·  Synthetic scenario prototype. Human review required.",
    )

    # ── Slide 2: Organisation Context ──────────────────────────────────────────
    ctx = []
    if org_name:
        ctx.append(f"Organisation: {org_name}")
    if org_type:
        ctx.append(f"Type: {org_type}")
    if sector:
        ctx.append(f"Sector: {sector}")
    if staff_count:
        ctx.append(f"Staff count: {staff_count}")
    if roles:
        ctx.append(f"Roles: {', '.join(str(r) for r in roles[:4])}")
    if current_ai_use:
        ctx.append(f"Current AI use: {format_slide_text(current_ai_use, 80)}")
    add_two_column_slide(
        prs,
        title="Organisation Context",
        left_title="Organisation Overview",
        left_bullets=ctx[:6] or ["No scenario loaded."],
        right_title="Main Concerns",
        right_bullets=extract_slide_bullets(concerns, 6) or ["—"],
    )

    # ── Slide 3: Training Need ──────────────────────────────────────────────────
    assessment  = pack_data.get("training_needs_assessment") or {}
    outcomes    = assessment.get("recommended_learning_outcomes") or []
    need_bullets = extract_slide_bullets(outcomes, 6)
    add_bullet_slide(
        prs,
        title="Training Need",
        bullets=need_bullets or ["Training Needs Assessment not yet generated."],
        speaker_notes="Learning outcomes derived from the scenario's priority topics and staff concerns.",
    )

    # ── Slide 4: Priority Topics ────────────────────────────────────────────────
    topic_assessments = assessment.get("topic_assessments") or []
    high_topics = [t.get("title", "") for t in topic_assessments if t.get("priority_level") == "high"]
    med_topics  = [t.get("title", "") for t in topic_assessments if t.get("priority_level") == "medium"]
    if not high_topics and not med_topics:
        high_topics = list(priority_topics[:6])
    add_two_column_slide(
        prs,
        title="Priority Training Topics",
        left_title="HIGH Priority",
        left_bullets=high_topics[:6] or ["Topics not yet scored."],
        right_title="MEDIUM Priority",
        right_bullets=med_topics[:6] or ["—"],
    )

    # ── Slide 5: Workshop Plan ──────────────────────────────────────────────────
    workshop_plan = pack_data.get("workshop_plan") or {}
    if workshop_plan:
        plan_left = [
            f"Title: {format_slide_text(workshop_plan.get('workshop_title', ''), 70)}",
            f"Duration: {workshop_plan.get('duration_minutes', '')} minutes",
            f"Delivery: {workshop_plan.get('delivery_mode', '')}",
        ]
        if roles:
            plan_left.append(f"Audience: {', '.join(str(r) for r in roles[:3])}")
        add_two_column_slide(
            prs,
            title="Workshop Plan",
            left_title="Session Details",
            left_bullets=plan_left,
            right_title="Responsible-Use Messages",
            right_bullets=extract_slide_bullets(
                workshop_plan.get("responsible_use_messages") or [], 6
            ) or ["—"],
        )
    else:
        add_bullet_slide(prs, "Workshop Plan", ["Workshop Plan not yet generated."])

    # ── Slide 6: Timed Agenda ───────────────────────────────────────────────────
    agenda = (workshop_plan.get("agenda") or []) if workshop_plan else []
    if agenda:
        agenda_bullets = [
            f"{item.get('time_range', '')}  —  {item.get('section_title', '')}"
            for item in agenda[:6]
        ]
        add_bullet_slide(
            prs,
            title="Workshop Agenda",
            bullets=agenda_bullets,
            speaker_notes="Full timing and trainer notes are in the workshop plan Markdown export.",
        )
    else:
        add_bullet_slide(prs, "Workshop Agenda", ["Workshop Plan not yet generated."])

    # ── Slide 7: Training Activities ────────────────────────────────────────────
    activities = pack_data.get("training_activities") or []
    if activities:
        act_bullets = [
            f"{a.get('activity_title', 'Activity')} ({a.get('duration_minutes', '')} min)"
            for a in activities[:6]
        ]
        type_counts = analytics.get("activity_type_counts", {})
        type_lines  = [f"{k.replace('_', ' ').title()}: {v}" for k, v in list(type_counts.items())[:6]]
        add_two_column_slide(
            prs,
            title="Training Activities",
            left_title=f"Activities  ({len(activities)} total)",
            left_bullets=act_bullets,
            right_title="Activity Mix",
            right_bullets=type_lines or ["Analytics not available."],
        )
    else:
        add_bullet_slide(prs, "Training Activities", ["Training Activities not yet generated."])

    # ── Slide 8: Safe AI Use Rules ──────────────────────────────────────────────
    handout  = pack_data.get("staff_handout") or {}
    allowed  = extract_slide_bullets(handout.get("allowed_ai_uses") or [], 6)
    add_bullet_slide(
        prs,
        title="What Staff CAN Use AI For",
        bullets=allowed or ["Staff Handout not yet generated."],
        speaker_notes="These are the approved, safe uses for this organisation's staff.",
    )

    # ── Slide 9: What Staff Must Not Do ─────────────────────────────────────────
    prohibited = extract_slide_bullets(handout.get("prohibited_ai_uses") or [], 6)
    add_bullet_slide(
        prs,
        title="What Staff Must NOT Do",
        bullets=prohibited or ["Staff Handout not yet generated."],
        speaker_notes="These are absolute prohibitions — not matters of individual discretion.",
    )

    # ── Slide 10: Human Review Checklist ────────────────────────────────────────
    hr_items = handout.get("human_review_checklist") or []
    checklist_bullets = [f"☐  {format_slide_text(str(c), 100)}" for c in hr_items[:6]]
    add_bullet_slide(
        prs,
        title="Human Review Checklist",
        bullets=checklist_bullets or (
            ["Staff Handout not yet generated."] if not handout else ["No checklist items available."]
        ),
        speaker_notes="Staff should complete this checklist before acting on any AI-generated output.",
    )

    # ── Slide 11: Escalation Guidance ───────────────────────────────────────────
    escalation = handout.get("escalation_guidance") or [] if handout else []
    esc_bullets = []
    if isinstance(escalation, list):
        esc_bullets = [
            f"{e.get('issue', '')}: {format_slide_text(e.get('what_to_do', ''), 80)}"
            for e in escalation[:6]
            if isinstance(e, dict)
        ]
    add_bullet_slide(
        prs,
        title="Escalation Guidance",
        bullets=esc_bullets or (
            ["Staff Handout not yet generated."] if not handout else ["No escalation data available."]
        ),
        speaker_notes="Staff must follow these routes — AI must never substitute for escalation.",
    )

    # ── Slide 12: Knowledge Check ────────────────────────────────────────────────
    kc = pack_data.get("knowledge_check") or {}
    if kc:
        mcqs        = kc.get("multiple_choice_questions") or []
        scenarios_q = kc.get("scenario_questions") or []
        reflections = kc.get("reflection_questions") or []
        kc_left = [
            f"Multiple-choice: {len(mcqs)} questions",
            f"Scenario questions: {len(scenarios_q)}",
            f"Reflection questions: {len(reflections)}",
            f"Answer key: {'Included' if kc.get('answer_key') else 'Staff copy — no answers'}",
            f"Pass guidance: {format_slide_text(kc.get('pass_guidance', ''), 80)}",
        ]
        kc_topics   = analytics.get("knowledge_check_topic_counts", {})
        topic_lines = [f"{k.replace('_', ' ').title()}: {v} Q" for k, v in list(kc_topics.items())[:6]]
        add_two_column_slide(
            prs,
            title="Knowledge Check",
            left_title="Check Overview",
            left_bullets=kc_left,
            right_title="Topics Covered",
            right_bullets=topic_lines or ["Topic analytics not available."],
        )
    else:
        add_bullet_slide(prs, "Knowledge Check", ["Knowledge Check not yet generated."])

    # ── Slide 13: Pack Completion ────────────────────────────────────────────────
    section_completion = analytics.get("section_completion", {})
    report_quality     = analytics.get("report_quality", {})
    if section_completion:
        completion_bullets = [
            f"{'✓' if done else '○'}  {sec}"
            for sec, done in section_completion.items()
        ]
        avail = report_quality.get("sections_available", "—")
        total = report_quality.get("sections_total", "—")
        pct   = report_quality.get("completeness_pct", "—")
        pct_str = f"{pct:.0f}%" if isinstance(pct, (int, float)) else str(pct)
        qa_bullets = [
            f"Sections available: {avail} of {total}",
            f"Completeness: {pct_str}",
            "Review all sections before delivery.",
            "Confirm all examples are synthetic.",
            "Get responsible-owner sign-off before use.",
        ]
        add_two_column_slide(
            prs,
            title="Pack Completion",
            left_title="Section Status",
            left_bullets=completion_bullets[:6],
            right_title="Quality Reminders",
            right_bullets=qa_bullets,
        )
    else:
        add_bullet_slide(
            prs,
            title="Pack Completion",
            bullets=["Generate the training pack first to see completion status."],
        )

    # ── Slide 14: Responsible-Use Boundaries ─────────────────────────────────────
    add_responsible_use_slide(prs)

    # ── Slide 15: Recommended Next Steps ─────────────────────────────────────────
    next_steps = [
        "Review all generated sections before use.",
        "Confirm all examples are synthetic — add organisation-specific context.",
        "Share the facilitator guide with the trainer one week before delivery.",
        "Confirm the safeguarding escalation route with the designated safeguarding lead.",
        "Run a pilot workshop with a small group before full delivery.",
        "Align materials with the organisation's actual policies before use.",
    ]
    add_bullet_slide(
        prs,
        title="Recommended Next Steps",
        bullets=next_steps,
        speaker_notes=f"Build 4 · AI Staff Training and Workshop Generator · {today}",
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
