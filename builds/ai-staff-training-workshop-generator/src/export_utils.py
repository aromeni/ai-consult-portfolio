"""Export utilities for the AI Staff Training and Workshop Generator.

Generates PDF and PowerPoint exports from training pack data.
All content is deterministic and based on synthetic scenario data only.
No external API calls. No real data.
"""

import io
import os
import re
import textwrap
from datetime import date


# ── Filename helpers ────────────────────────────────────────────────────────────

def create_safe_filename(title: str, extension: str) -> str:
    """Return a kebab-case filename safe for download."""
    if not isinstance(title, str):
        title = "training-pack"
    slug = re.sub(r"[^a-z0-9]+", "-", title.lower().strip()).strip("-")
    slug = slug[:60] or "training-pack"
    ext = extension.lstrip(".")
    return f"{slug}.{ext}"


# ── Text helpers ────────────────────────────────────────────────────────────────

def format_text_for_pdf(text: str) -> str:
    """Strip Markdown formatting for PDF paragraph use."""
    if not isinstance(text, str):
        return ""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"^[-*]\s+", "• ", text, flags=re.MULTILINE)
    return text.strip()


def format_slide_text(text: str, max_chars: int = 900) -> str:
    """Truncate text for slide bodies."""
    if not isinstance(text, str):
        return ""
    text = format_text_for_pdf(text)
    if len(text) > max_chars:
        text = text[:max_chars].rsplit(" ", 1)[0] + "…"
    return text


def _bullet_list(items: list, max_items: int = 8) -> str:
    """Format a list as bullet lines for slides/PDF."""
    if not items:
        return ""
    return "\n".join(f"• {str(i)}" for i in items[:max_items])


# ── Chart generation (matplotlib) ──────────────────────────────────────────────

def _chart_bytes(fig) -> bytes:
    """Save a matplotlib figure to PNG bytes and close it."""
    import matplotlib
    matplotlib.use("Agg")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=100, bbox_inches="tight")
    buf.seek(0)
    data = buf.read()
    buf.close()
    import matplotlib.pyplot as plt
    plt.close(fig)
    return data


def generate_chart_images_for_training_pack(
    pack_data: dict,
    analytics: dict,
    output_dir: str = "outputs/charts",
) -> dict:
    """Generate chart images and return {chart_name: bytes}.

    Returns an empty dict if matplotlib is unavailable or data is missing.
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    charts: dict = {}

    # A. Section completion bar chart
    section_completion = analytics.get("section_completion", {})
    if section_completion:
        labels = list(section_completion.keys())
        values = [1 if v else 0 for v in section_completion.values()]
        colors = ["#2ecc71" if v else "#bdc3c7" for v in values]
        fig, ax = plt.subplots(figsize=(7, 3.5))
        bars = ax.barh(labels, values, color=colors, height=0.5)
        ax.set_xlim(0, 1.2)
        ax.set_xticks([])
        ax.set_title("Section Completion", fontsize=12, fontweight="bold", pad=10)
        for bar, val in zip(bars, values):
            ax.text(
                bar.get_width() + 0.03,
                bar.get_y() + bar.get_height() / 2,
                "✓ Complete" if val else "○ Pending",
                va="center", fontsize=9,
                color="#27ae60" if val else "#7f8c8d",
            )
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["bottom"].set_visible(False)
        fig.tight_layout()
        charts["section_completion"] = _chart_bytes(fig)

    # B. Training topic coverage
    topic_counts = analytics.get("topic_counts", {})
    if topic_counts:
        labels = [str(k)[:28] for k in topic_counts.keys()]
        values = list(topic_counts.values())
        fig, ax = plt.subplots(figsize=(7, max(3, len(labels) * 0.45)))
        ax.barh(labels, values, color="#3498db", height=0.5)
        ax.set_xlabel("Count")
        ax.set_title("Priority Training Topics", fontsize=12, fontweight="bold", pad=10)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        fig.tight_layout()
        charts["topic_coverage"] = _chart_bytes(fig)

    # C. Activity mix pie chart
    activity_type_counts = analytics.get("activity_type_counts", {})
    if len(activity_type_counts) > 1:
        labels = [str(k).replace("_", " ").title()[:22] for k in activity_type_counts.keys()]
        values = list(activity_type_counts.values())
        palette = ["#3498db", "#2ecc71", "#e74c3c", "#f39c12", "#9b59b6", "#1abc9c", "#e67e22", "#34495e"]
        fig, ax = plt.subplots(figsize=(6, 4))
        wedges, texts, autotexts = ax.pie(
            values, labels=labels,
            colors=palette[:len(values)],
            autopct="%1.0f%%", startangle=90,
            textprops={"fontsize": 9},
        )
        ax.set_title("Activity Mix", fontsize=12, fontweight="bold", pad=10)
        fig.tight_layout()
        charts["activity_mix"] = _chart_bytes(fig)
    elif len(activity_type_counts) == 1:
        # Single type — skip pie, just note it
        pass

    # D. Workshop time allocation
    time_allocation = analytics.get("workshop_time_allocation", [])
    time_allocation = [t for t in time_allocation if isinstance(t, dict) and t.get("minutes", 0) > 0]
    if time_allocation:
        labels = [str(t["section"])[:24] for t in time_allocation]
        values = [int(t["minutes"]) for t in time_allocation]
        fig, ax = plt.subplots(figsize=(7, max(3, len(labels) * 0.45)))
        ax.barh(labels, values, color="#9b59b6", height=0.5)
        ax.set_xlabel("Minutes")
        ax.set_title("Workshop Time Allocation", fontsize=12, fontweight="bold", pad=10)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        fig.tight_layout()
        charts["time_allocation"] = _chart_bytes(fig)

    # E. Knowledge check topic coverage
    kc_topics = analytics.get("knowledge_check_topic_counts", {})
    if kc_topics:
        labels = [str(k).replace("_", " ").title()[:24] for k in kc_topics.keys()]
        values = list(kc_topics.values())
        fig, ax = plt.subplots(figsize=(7, max(3, len(labels) * 0.45)))
        ax.barh(labels, values, color="#e67e22", height=0.5)
        ax.set_xlabel("Questions")
        ax.set_title("Knowledge Check Topic Coverage", fontsize=12, fontweight="bold", pad=10)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        fig.tight_layout()
        charts["kc_topic_coverage"] = _chart_bytes(fig)

    return charts


# ── PDF export (reportlab) ──────────────────────────────────────────────────────

def export_training_pack_to_pdf(
    pack_data: dict,
    markdown_text: str,
    analytics: dict | None = None,
    output_path: str | None = None,
) -> bytes:
    """Generate a PDF training pack and return its bytes.

    Suitable for st.download_button(data=bytes, mime="application/pdf").
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer,
        HRFlowable, PageBreak, Image as RLImage,
    )

    if not isinstance(pack_data, dict):
        pack_data = {}
    if analytics is None:
        analytics = {}

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2.5 * cm,
        bottomMargin=2.5 * cm,
        title="AI Staff Training Pack",
    )

    styles = getSampleStyleSheet()
    W = A4[0] - 4 * cm  # usable width

    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=18, spaceAfter=12,
                        textColor=colors.HexColor("#1a2744"), leading=22)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=14, spaceAfter=8,
                        spaceBefore=14, textColor=colors.HexColor("#2c3e50"), leading=18)
    h3 = ParagraphStyle("H3", parent=styles["Heading3"], fontSize=11, spaceAfter=6,
                        spaceBefore=10, textColor=colors.HexColor("#34495e"))
    body = ParagraphStyle("Body", parent=styles["Normal"], fontSize=10, leading=14,
                          spaceAfter=6)
    small = ParagraphStyle("Small", parent=styles["Normal"], fontSize=9, leading=12,
                            textColor=colors.HexColor("#7f8c8d"), spaceAfter=4)
    notice = ParagraphStyle("Notice", parent=styles["Normal"], fontSize=9, leading=12,
                             textColor=colors.HexColor("#e67e22"), spaceAfter=6,
                             borderPad=4)
    bullet = ParagraphStyle("Bullet", parent=styles["Normal"], fontSize=10, leading=13,
                             leftIndent=12, spaceAfter=3)

    scenario = pack_data.get("scenario") or {}
    org_name = scenario.get("organisation_name", "Organisation")
    org_type = scenario.get("organisation_type", "")
    sector = scenario.get("sector", "")
    staff_count = scenario.get("staff_count", "")
    today = str(date.today())

    report_quality = analytics.get("report_quality", {})
    sections_available = report_quality.get("sections_available", 0)
    sections_total = report_quality.get("sections_total", 7)

    story = []

    # ── Cover page ──────────────────────────────────────────────────────────────
    story.append(Spacer(1, 3 * cm))
    story.append(Paragraph("Responsible AI Staff Training Pack", h1))
    story.append(HRFlowable(width=W, thickness=2, color=colors.HexColor("#3498db")))
    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph(org_name, ParagraphStyle("OrgName", parent=h2, fontSize=16)))
    if org_type:
        story.append(Paragraph(org_type, small))
    if sector:
        story.append(Paragraph(f"Sector: {sector}", small))
    if staff_count:
        story.append(Paragraph(f"Staff: {staff_count}", small))
    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph(f"Generated: {today}", small))
    story.append(Spacer(1, 1 * cm))
    story.append(Paragraph(
        "⚠ PROTOTYPE — All scenarios are synthetic. Outputs require human review "
        "before use in a real training session. This document does not constitute "
        "legal, safeguarding, HR, compliance, medical, or professional advice.",
        notice,
    ))
    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph(
        f"Sections included: {sections_available} of {sections_total} generated outputs.",
        small,
    ))
    story.append(PageBreak())

    # ── Executive summary ───────────────────────────────────────────────────────
    story.append(Paragraph("Executive Summary", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    _wp = pack_data.get("workshop_plan") or {}
    exec_rows = [
        ("Organisation", org_name),
        ("Training goal", format_text_for_pdf(scenario.get("training_goal", ""))),
        ("Audience", ", ".join(str(r) for r in (scenario.get("staff_roles") or []))),
        ("Duration", f"{_wp.get('duration_minutes', '')} minutes" if _wp.get("duration_minutes") else scenario.get("training_duration", "")),
        ("Delivery mode", _wp.get("delivery_mode", "") or scenario.get("delivery_mode", "")),
        ("Sections included", f"{sections_available} of {sections_total} generated outputs"),
    ]
    for label, value in exec_rows:
        if value:
            story.append(Paragraph(f"<b>{label}:</b> {value}", body))
    story.append(Spacer(1, 0.4 * cm))

    # ── Section completion ──────────────────────────────────────────────────────
    charts = {}
    try:
        charts = generate_chart_images_for_training_pack(pack_data, analytics)
    except Exception:
        pass

    story.append(Paragraph("Training Pack Overview", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))

    section_completion = analytics.get("section_completion", {})
    for section, complete in section_completion.items():
        mark = "✓" if complete else "○"
        colour = "#27ae60" if complete else "#95a5a6"
        story.append(Paragraph(
            f"<font color='{colour}'>{mark}</font>  {section}",
            bullet,
        ))

    if "section_completion" in charts:
        try:
            img_buf = io.BytesIO(charts["section_completion"])
            img = RLImage(img_buf, width=W, height=W * 0.45)
            story.append(Spacer(1, 0.3 * cm))
            story.append(img)
        except Exception:
            pass

    story.append(PageBreak())

    # ── Training context ────────────────────────────────────────────────────────
    story.append(Paragraph("1. Organisation Scenario", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    if scenario:
        story.append(Paragraph(f"<b>Organisation:</b> {org_name}", body))
        if org_type:
            story.append(Paragraph(f"<b>Type:</b> {org_type}", body))
        if sector:
            story.append(Paragraph(f"<b>Sector:</b> {sector}", body))
        if staff_count:
            story.append(Paragraph(f"<b>Staff count:</b> {staff_count}", body))
        training_goal = format_text_for_pdf(scenario.get("training_goal", ""))
        if training_goal:
            story.append(Paragraph(f"<b>Training goal:</b> {training_goal}", body))
        current_ai_use = format_text_for_pdf(scenario.get("current_ai_use", ""))
        if current_ai_use:
            story.append(Paragraph(f"<b>Current AI use:</b> {current_ai_use}", body))
        roles = scenario.get("staff_roles") or []
        if roles:
            story.append(Paragraph(f"<b>Staff roles:</b> {', '.join(str(r) for r in roles)}", body))
        topics = scenario.get("priority_topics") or []
        if topics:
            story.append(Paragraph("<b>Priority topics:</b>", body))
            for t in topics:
                story.append(Paragraph(f"• {t}", bullet))
    else:
        story.append(Paragraph("No scenario data available.", small))
    story.append(PageBreak())

    # ── Training needs ──────────────────────────────────────────────────────────
    story.append(Paragraph("2. Training Needs Assessment", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    assessment = pack_data.get("training_needs_assessment") or {}
    if assessment:
        topic_assessments = assessment.get("topic_assessments") or []
        if topic_assessments:
            story.append(Paragraph("<b>Topic Priorities:</b>", body))
            for t in topic_assessments[:10]:
                level = t.get("priority_level", "").upper()
                title = t.get("title", "")
                story.append(Paragraph(f"• [{level}] {title}", bullet))
        outcomes = assessment.get("recommended_learning_outcomes") or []
        if outcomes:
            story.append(Spacer(1, 0.3 * cm))
            story.append(Paragraph("<b>Learning Outcomes:</b>", body))
            for o in outcomes:
                story.append(Paragraph(f"• {format_text_for_pdf(o)}", bullet))
    else:
        story.append(Paragraph("Training Needs Assessment not generated.", small))

    if "topic_coverage" in charts:
        try:
            img_buf = io.BytesIO(charts["topic_coverage"])
            img = RLImage(img_buf, width=W, height=W * 0.5)
            story.append(Spacer(1, 0.3 * cm))
            story.append(img)
        except Exception:
            pass
    story.append(PageBreak())

    # ── Workshop plan ────────────────────────────────────────────────────────────
    story.append(Paragraph("3. Workshop Plan", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    workshop_plan = pack_data.get("workshop_plan") or {}
    if workshop_plan:
        story.append(Paragraph(f"<b>Title:</b> {format_text_for_pdf(workshop_plan.get('workshop_title', ''))}", body))
        story.append(Paragraph(f"<b>Duration:</b> {workshop_plan.get('duration_minutes', '')} minutes", body))
        story.append(Paragraph(f"<b>Delivery:</b> {workshop_plan.get('delivery_mode', '')}", body))
        agenda = workshop_plan.get("agenda") or []
        if agenda:
            story.append(Paragraph("<b>Agenda:</b>", body))
            for item in agenda:
                time_range = item.get("time_range", "")
                title = item.get("section_title", "")
                story.append(Paragraph(f"• {time_range} — {title}", bullet))
    else:
        story.append(Paragraph("Workshop Plan not generated.", small))

    if "time_allocation" in charts:
        try:
            img_buf = io.BytesIO(charts["time_allocation"])
            img = RLImage(img_buf, width=W, height=W * 0.45)
            story.append(Spacer(1, 0.3 * cm))
            story.append(img)
        except Exception:
            pass
    story.append(PageBreak())

    # ── Activities ───────────────────────────────────────────────────────────────
    story.append(Paragraph("4. Training Activities", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    activities = pack_data.get("training_activities") or []
    if activities:
        story.append(Paragraph(f"<b>Total activities:</b> {len(activities)}", body))
        for i, act in enumerate(activities[:8], 1):
            title = act.get("activity_title", "Activity")
            duration = act.get("duration_minutes", "")
            objective = format_text_for_pdf(act.get("learning_objective", ""))
            story.append(Paragraph(f"<b>{i}. {title}</b> ({duration} min)", h3))
            if objective:
                story.append(Paragraph(objective, body))
    else:
        story.append(Paragraph("Training Activities not generated.", small))

    if "activity_mix" in charts:
        try:
            img_buf = io.BytesIO(charts["activity_mix"])
            img = RLImage(img_buf, width=W * 0.7, height=W * 0.5)
            story.append(Spacer(1, 0.3 * cm))
            story.append(img)
        except Exception:
            pass
    story.append(PageBreak())

    # ── Facilitator guide summary ────────────────────────────────────────────────
    story.append(Paragraph("5. Facilitator Guide", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    guide = pack_data.get("facilitator_guide") or {}
    if guide:
        if guide.get("session_purpose"):
            story.append(Paragraph(
                f"<b>Session purpose:</b> {format_text_for_pdf(guide.get('session_purpose', ''))}",
                body,
            ))
        principles = guide.get("facilitator_principles") or []
        if principles:
            story.append(Paragraph("<b>Facilitator Principles:</b>", body))
            for p in principles[:6]:
                story.append(Paragraph(f"• {format_text_for_pdf(p)}", bullet))
        prep = guide.get("preparation_checklist") or []
        if prep:
            story.append(Paragraph("<b>Preparation Checklist:</b>", body))
            for i, item in enumerate(prep[:8], 1):
                story.append(Paragraph(f"{i}. {format_text_for_pdf(item)}", bullet))
        misconceptions = guide.get("common_misconceptions") or []
        if misconceptions:
            story.append(Paragraph("<b>Common Misconceptions to Address:</b>", body))
            for m in misconceptions[:5]:
                story.append(Paragraph(
                    f"• {format_text_for_pdf(m.get('misconception', '') if isinstance(m, dict) else str(m))}",
                    bullet,
                ))
    else:
        story.append(Paragraph("Facilitator Guide not generated.", small))
    story.append(PageBreak())

    # ── Staff handout summary ────────────────────────────────────────────────────
    story.append(Paragraph("6. Staff Handout — Safe-Use Reference", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    handout = pack_data.get("staff_handout") or {}
    if handout:
        allowed = handout.get("allowed_ai_uses") or []
        prohibited = handout.get("prohibited_ai_uses") or []
        principles = handout.get("safe_use_principles") or []
        if principles:
            story.append(Paragraph("<b>Safe-Use Principles (summary):</b>", body))
            for p in principles[:5]:
                story.append(Paragraph(f"• {format_text_for_pdf(p)}", bullet))
        if allowed:
            story.append(Paragraph("<b>Staff Can Use AI For:</b>", body))
            for a in allowed[:5]:
                story.append(Paragraph(f"• {format_text_for_pdf(a)}", bullet))
        if prohibited:
            story.append(Paragraph("<b>Staff Must NOT Use AI For:</b>", body))
            for p in prohibited[:5]:
                story.append(Paragraph(f"• {format_text_for_pdf(p)}", bullet))
    else:
        story.append(Paragraph("Staff Handout not generated.", small))
    story.append(PageBreak())

    # ── Knowledge check summary ──────────────────────────────────────────────────
    story.append(Paragraph("7. Knowledge Check", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    kc = pack_data.get("knowledge_check") or {}
    if kc:
        mcqs = kc.get("multiple_choice_questions") or []
        scenarios_q = kc.get("scenario_questions") or []
        reflections = kc.get("reflection_questions") or []
        story.append(Paragraph(f"<b>Multiple-choice questions:</b> {len(mcqs)}", body))
        story.append(Paragraph(f"<b>Scenario questions:</b> {len(scenarios_q)}", body))
        story.append(Paragraph(f"<b>Reflection questions:</b> {len(reflections)}", body))
        story.append(Paragraph(f"<b>Answer key:</b> {'Included' if kc.get('answer_key') else 'Not included'}", body))
        if mcqs:
            story.append(Spacer(1, 0.3 * cm))
            story.append(Paragraph("<b>Sample MCQs:</b>", body))
            for q in mcqs[:3]:
                story.append(Paragraph(f"• {format_text_for_pdf(q.get('question', ''))}", bullet))
    else:
        story.append(Paragraph("Knowledge Check not generated.", small))

    if "kc_topic_coverage" in charts:
        try:
            img_buf = io.BytesIO(charts["kc_topic_coverage"])
            img = RLImage(img_buf, width=W, height=W * 0.4)
            story.append(Spacer(1, 0.3 * cm))
            story.append(img)
        except Exception:
            pass
    story.append(PageBreak())

    # ── Facilitator review checklist ──────────────────────────────────────────────
    story.append(Paragraph("8. Facilitator Review Checklist", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    story.append(Paragraph("Before delivering this training pack, confirm:", body))
    review_checklist = [
        "All sections have been read in full by the facilitator.",
        "All examples and scenario cards are confirmed synthetic.",
        "Allowed and prohibited AI uses match the organisation's actual policy.",
        "The safeguarding escalation route names the correct contact.",
        "The approved tools list reflects what the organisation has approved.",
        "The answer key has been removed from the staff copy.",
        "Timings have been checked against the booked session length.",
        "A responsible owner has approved the materials for delivery.",
    ]
    for item in review_checklist:
        story.append(Paragraph(f"☐  {item}", bullet))
    story.append(PageBreak())

    # ── Responsible-use boundaries ────────────────────────────────────────────────
    story.append(Paragraph("9. Responsible-Use Boundaries", h2))
    story.append(HRFlowable(width=W, thickness=2, color=colors.HexColor("#e74c3c")))
    responsible_use_items = [
        "Use synthetic or approved scenarios only — no real learner, safeguarding, HR, personal, or regulated data.",
        "No external AI API calls — all generation runs locally and deterministically.",
        "Human review is required before using any output in a real training session.",
        "This pack does not constitute legal, safeguarding, HR, compliance, medical, or professional advice.",
        "Training materials are support materials, not official organisational policy.",
        "Production use requires governance, DPIA, security, privacy, and responsible-owner approval.",
    ]
    for item in responsible_use_items:
        story.append(Paragraph(f"• {item}", bullet))
    story.append(PageBreak())

    # ── Prototype limitations ─────────────────────────────────────────────────────
    story.append(Paragraph("10. Prototype Limitations", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    limitations = [
        "This is a production-style prototype, not a certified training compliance system.",
        "Session state is in-memory only — resets on app restart.",
        "All scenarios are synthetic — not based on any real organisation.",
        "No authentication or access control.",
        "No staff completion tracking or LMS integration.",
        "Outputs require qualified human review before use in real training delivery.",
    ]
    for item in limitations:
        story.append(Paragraph(f"• {item}", bullet))
    story.append(PageBreak())

    # ── Recommended next steps ────────────────────────────────────────────────────
    story.append(Paragraph("11. Recommended Next Steps", h2))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    next_steps = [
        "Review all generated sections before use.",
        "Confirm all examples are synthetic — add organisation-specific context as needed.",
        "Share the facilitator guide with the trainer at least one week before delivery.",
        "Confirm the safeguarding escalation route with the designated safeguarding lead.",
        "Confirm which AI tools are approved for use in this organisation.",
        "Run a pilot workshop with a small group before full delivery.",
        "Collect staff feedback after the session and update the training pack.",
        "Review the knowledge check results and follow up with staff who need support.",
        "Treat this pack as a starting point — align it with your organisation's actual policies.",
    ]
    for i, step in enumerate(next_steps, 1):
        story.append(Paragraph(f"{i}. {step}", bullet))

    story.append(Spacer(1, 1 * cm))
    story.append(HRFlowable(width=W, thickness=1, color=colors.HexColor("#bdc3c7")))
    story.append(Paragraph(
        f"Build 4 · AI Staff Training and Workshop Generator · BrightPath ChatGPT Mastery Project  "
        f"| Generated: {today}  | Synthetic scenario prototype. Human review required.",
        small,
    ))

    doc.build(story)
    buf.seek(0)
    return buf.read()


# ── PPTX export (python-pptx) ───────────────────────────────────────────────────

def export_training_pack_to_pptx(
    pack_data: dict,
    analytics: dict | None = None,
    output_path: str | None = None,
) -> bytes:
    """Generate a PPTX training pack presentation and return its bytes.

    Suitable for st.download_button(data=bytes, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation").
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not isinstance(pack_data, dict):
        pack_data = {}
    if analytics is None:
        analytics = {}

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colour palette
    NAVY = RGBColor(0x1A, 0x27, 0x44)
    BLUE = RGBColor(0x34, 0x98, 0xDB)
    DARK = RGBColor(0x2C, 0x3E, 0x50)
    GREY = RGBColor(0x7F, 0x8C, 0x8D)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    AMBER = RGBColor(0xE6, 0x7E, 0x22)
    GREEN = RGBColor(0x27, 0xAE, 0x60)

    scenario = pack_data.get("scenario") or {}
    org_name = scenario.get("organisation_name", "Organisation")
    org_type = scenario.get("organisation_type", "")
    sector = scenario.get("sector", "")
    staff_count = scenario.get("staff_count", "")
    roles = scenario.get("staff_roles") or []
    concerns = scenario.get("main_concerns") or []
    current_ai_use = scenario.get("current_ai_use", "")
    priority_topics = scenario.get("priority_topics") or []
    today = str(date.today())

    blank_layout = prs.slide_layouts[6]  # blank

    def add_slide():
        return prs.slides.add_slide(blank_layout)

    def add_textbox(slide, text, left, top, width, height,
                    font_size=20, bold=False, color=None, wrap=True, align=PP_ALIGN.LEFT):
        from pptx.util import Inches, Pt, Emu
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def add_title_bar(slide, title, subtitle=""):
        # Navy bar at top
        from pptx.util import Inches, Pt
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.33), Inches(1.4),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        add_textbox(slide, title, 0.3, 0.1, 12.5, 0.8,
                    font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        if subtitle:
            add_textbox(slide, subtitle, 0.3, 0.9, 12.5, 0.4,
                        font_size=13, bold=False, color=RGBColor(0xBD, 0xC3, 0xC7),
                        align=PP_ALIGN.LEFT)

    def add_bullets(slide, items, left, top, width, height,
                    font_size=15, color=None, max_items=10):
        if not items:
            return
        text = "\n".join(f"• {format_slide_text(str(i), 120)}" for i in items[:max_items])
        add_textbox(slide, text, left, top, width, height,
                    font_size=font_size, color=color or DARK)

    def add_notice(slide, text="Synthetic scenario prototype. Human review required.", top=6.8):
        add_textbox(slide, text, 0.3, top, 12.7, 0.4,
                    font_size=9, color=GREY, align=PP_ALIGN.LEFT)

    # ── Slide 1: Title ──────────────────────────────────────────────────────────
    slide = add_slide()
    # Full background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    add_textbox(slide, "Responsible AI Staff Training Pack",
                0.4, 1.8, 12.5, 1.2, font_size=36, bold=True, color=WHITE)
    add_textbox(slide, org_name, 0.4, 3.2, 12.5, 0.8,
                font_size=22, bold=False, color=BLUE)
    details = []
    if org_type:
        details.append(org_type)
    if sector:
        details.append(sector)
    if staff_count:
        details.append(f"{staff_count} staff")
    if details:
        add_textbox(slide, "  |  ".join(details), 0.4, 3.9, 12.5, 0.5,
                    font_size=13, color=RGBColor(0xBD, 0xC3, 0xC7))
    add_textbox(slide, f"Generated: {today}", 0.4, 4.6, 12.5, 0.4,
                font_size=12, color=GREY)
    add_textbox(slide,
                "⚠ Prototype — Synthetic scenarios only. Outputs require human review before use.",
                0.4, 6.6, 12.5, 0.6, font_size=10, color=RGBColor(0xE6, 0x7E, 0x22))

    # ── Slide 2: Training Context ───────────────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Training Context", "Organisation overview and AI use")
    col_items = []
    if org_name:
        col_items.append(f"Organisation: {org_name}")
    if org_type:
        col_items.append(f"Type: {org_type}")
    if sector:
        col_items.append(f"Sector: {sector}")
    if staff_count:
        col_items.append(f"Staff: {staff_count}")
    if roles:
        col_items.append(f"Roles: {', '.join(str(r) for r in roles[:4])}")
    if current_ai_use:
        col_items.append(f"Current AI use: {format_slide_text(current_ai_use, 80)}")
    add_bullets(slide, col_items, 0.4, 1.6, 6.0, 5.0, font_size=14)
    if concerns:
        add_textbox(slide, "Main Concerns:", 6.6, 1.6, 6.3, 0.4,
                    font_size=14, bold=True, color=DARK)
        add_bullets(slide, concerns, 6.6, 2.1, 6.3, 4.5, font_size=13, max_items=8)
    add_notice(slide, "All scenario data is synthetic — no real learner, safeguarding, HR, or personal data.")

    # ── Slide 3: Training Needs ─────────────────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Training Needs Assessment", "Priority topics and risk themes")
    assessment = pack_data.get("training_needs_assessment") or {}
    topic_assessments = assessment.get("topic_assessments") or []
    outcomes = assessment.get("recommended_learning_outcomes") or []
    high_topics = [t.get("title", "") for t in topic_assessments if t.get("priority_level") == "high"]
    med_topics = [t.get("title", "") for t in topic_assessments if t.get("priority_level") == "medium"]
    if high_topics:
        add_textbox(slide, "HIGH Priority:", 0.4, 1.6, 6.0, 0.4,
                    font_size=13, bold=True, color=RGBColor(0xE7, 0x4C, 0x3C))
        add_bullets(slide, high_topics, 0.4, 2.05, 6.0, 2.5, font_size=13, max_items=6)
    if med_topics:
        add_textbox(slide, "MEDIUM Priority:", 0.4, 4.0, 6.0, 0.4,
                    font_size=13, bold=True, color=RGBColor(0xF3, 0x9C, 0x12))
        add_bullets(slide, med_topics, 0.4, 4.45, 6.0, 1.8, font_size=12, max_items=4)
    if outcomes:
        add_textbox(slide, "Learning Outcomes:", 6.6, 1.6, 6.3, 0.4,
                    font_size=13, bold=True, color=DARK)
        add_bullets(slide, outcomes, 6.6, 2.05, 6.3, 4.8, font_size=12, max_items=6)
    if not assessment:
        add_textbox(slide, "Training Needs Assessment not yet generated.",
                    0.4, 2.0, 12.5, 0.6, font_size=14, color=GREY)
    add_notice(slide)

    # ── Slide 4: Workshop Plan ──────────────────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Workshop Plan", "Duration, delivery mode, and learning outcomes")
    workshop_plan = pack_data.get("workshop_plan") or {}
    if workshop_plan:
        details = [
            f"Title: {format_slide_text(workshop_plan.get('workshop_title', ''), 70)}",
            f"Duration: {workshop_plan.get('duration_minutes', '')} minutes",
            f"Delivery: {workshop_plan.get('delivery_mode', '')}",
            f"Audience: {', '.join(str(r) for r in (workshop_plan.get('audience') or [])[:3])}",
        ]
        add_bullets(slide, details, 0.4, 1.6, 6.0, 2.5, font_size=14)
        outcomes = workshop_plan.get("learning_outcomes") or []
        if outcomes:
            add_textbox(slide, "Learning Outcomes:", 0.4, 4.0, 6.0, 0.4,
                        font_size=13, bold=True, color=DARK)
            add_bullets(slide, outcomes, 0.4, 4.45, 6.0, 2.4, font_size=12, max_items=5)
        responsible_msgs = workshop_plan.get("responsible_use_messages") or []
        if responsible_msgs:
            add_textbox(slide, "Responsible-Use Messages to Cover:", 6.6, 1.6, 6.3, 0.4,
                        font_size=13, bold=True, color=DARK)
            add_bullets(slide, responsible_msgs, 6.6, 2.05, 6.3, 4.8, font_size=12, max_items=6)
    else:
        add_textbox(slide, "Workshop Plan not yet generated.",
                    0.4, 2.0, 12.5, 0.6, font_size=14, color=GREY)
    add_notice(slide)

    # ── Slide 5: Timed Agenda ───────────────────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Workshop Agenda", "Timed session plan")
    agenda = (workshop_plan.get("agenda") or []) if workshop_plan else []
    if agenda:
        y = 1.65
        for item in agenda[:8]:
            time_range = item.get("time_range", "")
            section_title = item.get("section_title", "")
            purpose = format_slide_text(item.get("purpose", ""), 60)
            line = f"{time_range}  —  {section_title}"
            add_textbox(slide, line, 0.4, y, 12.5, 0.35,
                        font_size=13, bold=False, color=DARK)
            if purpose:
                add_textbox(slide, purpose, 0.8, y + 0.35, 12.1, 0.3,
                            font_size=10, color=GREY)
                y += 0.75
            else:
                y += 0.45
            if y > 6.8:
                break
    else:
        add_textbox(slide, "Workshop Plan not yet generated.",
                    0.4, 2.0, 12.5, 0.6, font_size=14, color=GREY)
    add_notice(slide)

    # ── Slide 6: Training Activities ────────────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Training Activities", "Hands-on responsible AI practice")
    activities = pack_data.get("training_activities") or []
    if activities:
        add_textbox(slide, f"Total activities: {len(activities)}", 0.4, 1.6, 5.0, 0.4,
                    font_size=14, bold=True, color=DARK)
        activity_names = [
            f"{a.get('activity_title', 'Activity')} ({a.get('duration_minutes', '')} min)"
            for a in activities
        ]
        add_bullets(slide, activity_names, 0.4, 2.1, 6.0, 4.8, font_size=13, max_items=8)
        activity_type_counts = analytics.get("activity_type_counts", {})
        if activity_type_counts:
            add_textbox(slide, "Activity Types:", 6.6, 1.6, 6.3, 0.4,
                        font_size=13, bold=True, color=DARK)
            type_lines = [f"{k.replace('_', ' ').title()}: {v}" for k, v in activity_type_counts.items()]
            add_bullets(slide, type_lines, 6.6, 2.05, 6.3, 4.8, font_size=13, max_items=8)
    else:
        add_textbox(slide, "Training Activities not yet generated.",
                    0.4, 2.0, 12.5, 0.6, font_size=14, color=GREY)
    add_notice(slide, "All activity scenarios are synthetic — no real learner, safeguarding, or HR data.")

    # ── Slide 7: Safe-Use Rules ─────────────────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "AI Safe-Use Rules for Staff", "What staff can and cannot do")
    handout = pack_data.get("staff_handout") or {}
    allowed = handout.get("allowed_ai_uses") or []
    prohibited = handout.get("prohibited_ai_uses") or []
    if allowed:
        add_textbox(slide, "✓  Staff CAN use AI for:", 0.4, 1.6, 6.0, 0.45,
                    font_size=14, bold=True, color=GREEN)
        add_bullets(slide, allowed, 0.4, 2.1, 6.0, 4.8, font_size=13, max_items=6)
    if prohibited:
        add_textbox(slide, "✗  Staff MUST NOT use AI for:", 6.6, 1.6, 6.3, 0.45,
                    font_size=14, bold=True, color=RGBColor(0xE7, 0x4C, 0x3C))
        add_bullets(slide, prohibited, 6.6, 2.1, 6.3, 4.8, font_size=13, max_items=6)
    if not handout:
        add_textbox(slide, "Staff Handout not yet generated.",
                    0.4, 2.0, 12.5, 0.6, font_size=14, color=GREY)
    add_notice(slide)

    # ── Slide 8: Human Review and Escalation ────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Human Review and Escalation", "Before using AI output and who to contact")
    checklist = handout.get("human_review_checklist") or [] if handout else []
    escalation = handout.get("escalation_guidance") or [] if handout else []
    if checklist:
        add_textbox(slide, "Before using any AI output, check:", 0.4, 1.6, 6.0, 0.45,
                    font_size=13, bold=True, color=DARK)
        add_bullets(slide, [f"☐  {format_slide_text(c, 90)}" for c in checklist],
                    0.4, 2.1, 6.0, 4.8, font_size=12, max_items=7)
    if escalation:
        add_textbox(slide, "Escalation Guidance:", 6.6, 1.6, 6.3, 0.45,
                    font_size=13, bold=True, color=DARK)
        esc_lines = [
            f"{e.get('issue', '')}: {format_slide_text(e.get('what_to_do', ''), 70)}"
            for e in escalation[:5]
        ]
        add_bullets(slide, esc_lines, 6.6, 2.1, 6.3, 4.8, font_size=12, max_items=5)
    if not handout:
        add_textbox(slide, "Staff Handout not yet generated.",
                    0.4, 2.0, 12.5, 0.6, font_size=14, color=GREY)
    add_notice(slide)

    # ── Slide 9: Knowledge Check ────────────────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Knowledge Check", "MCQs, scenarios, and reflections")
    kc = pack_data.get("knowledge_check") or {}
    if kc:
        mcqs = kc.get("multiple_choice_questions") or []
        scenarios_q = kc.get("scenario_questions") or []
        reflections = kc.get("reflection_questions") or []
        stats = [
            f"Multiple-choice questions: {len(mcqs)}",
            f"Scenario questions: {len(scenarios_q)}",
            f"Reflection questions: {len(reflections)}",
            f"Answer key: {'Included' if kc.get('answer_key') else 'Not included'}",
            f"Pass guidance: {format_slide_text(kc.get('pass_guidance', ''), 80)}",
        ]
        add_bullets(slide, stats, 0.4, 1.6, 6.0, 3.0, font_size=14)
        kc_topics = analytics.get("knowledge_check_topic_counts", {})
        if kc_topics:
            add_textbox(slide, "Topics covered:", 6.6, 1.6, 6.3, 0.45,
                        font_size=13, bold=True, color=DARK)
            topic_lines = [f"{k.replace('_', ' ').title()}: {v} question(s)"
                           for k, v in list(kc_topics.items())[:8]]
            add_bullets(slide, topic_lines, 6.6, 2.1, 6.3, 4.8, font_size=13, max_items=8)
        if mcqs:
            add_textbox(slide, "Sample question:", 0.4, 4.8, 12.5, 0.35,
                        font_size=12, bold=True, color=DARK)
            add_textbox(slide, format_slide_text(mcqs[0].get("question", ""), 120),
                        0.4, 5.2, 12.5, 0.5, font_size=12, color=DARK)
    else:
        add_textbox(slide, "Knowledge Check not yet generated.",
                    0.4, 2.0, 12.5, 0.6, font_size=14, color=GREY)
    add_notice(slide)

    # ── Slide 10: Responsible-Use Boundaries ────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Responsible-Use Boundaries", "Non-negotiable constraints")
    ru_items = [
        "Use synthetic or approved scenarios only — no real learner, safeguarding, HR, personal, or regulated data.",
        "No external AI API calls — all generation is local and deterministic.",
        "Human review required before using any output in a real training session.",
        "This pack does not constitute legal, safeguarding, HR, compliance, medical, or professional advice.",
        "Training materials are support materials — not official organisational policy.",
        "Production use requires governance, DPIA, security, privacy, and responsible-owner approval.",
    ]
    add_bullets(slide, ru_items, 0.4, 1.6, 12.5, 5.5, font_size=14, max_items=6)
    # Red border emphasis
    border = slide.shapes.add_shape(1, Inches(0.2), Inches(1.5), Inches(0.1), Inches(5.5))
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
    border.line.fill.background()
    add_notice(slide)

    # ── Slide 11: Recommended Next Steps ────────────────────────────────────────
    slide = add_slide()
    add_title_bar(slide, "Recommended Next Steps", "Before and after delivery")
    next_steps = [
        "Review all generated sections before use.",
        "Confirm all examples are synthetic — add organisation-specific context.",
        "Share the facilitator guide with the trainer one week before delivery.",
        "Confirm the safeguarding escalation route with the designated safeguarding lead.",
        "Confirm which AI tools are approved for this organisation.",
        "Run a pilot workshop with a small group before full delivery.",
        "Collect staff feedback and update the training pack.",
        "Follow up with staff who need additional support after the knowledge check.",
        "Align training materials with the organisation's actual policies before use.",
    ]
    add_bullets(slide, next_steps, 0.4, 1.6, 12.5, 5.5, font_size=14, max_items=9)
    add_notice(slide,
               f"Build 4 · AI Staff Training and Workshop Generator · {today} · Synthetic prototype. Human review required.")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
